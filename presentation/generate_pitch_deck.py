"""
Generate a 5-slide pitch deck for Smart Financial Twin.

Output: presentation/slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
IDBI_BLUE = RGBColor(0, 102, 204)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Slide 1: Title Slide ---
def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Smart Financial Twin"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = IDBI_BLUE
    
    subtitle.text = "AI-Powered Personal Finance Advisor\nIDBI InnovaTech 2026\nTeam: [Your Team Name]\nJuly 2026"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = BLACK

# --- Slide 2: Problem Statement ---
def add_problem_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, content = slide.shapes.title, slide.placeholders[1]
    
    title.text = "The Gap in Personal Financial Planning"
    title.text_frame.paragraphs[0].font.color.rgb = IDBI_BLUE
    
    tf = content.text_frame
    tf.clear()
    
    problems = [
        "60% of Indians lack access to personalized financial advice.",
        "Traditional banking apps provide static insights (no 'what-if' scenarios).",
        "Fraud detection is reactive, not predictive.",
        "No unified view of financial health (savings, loans, investments)."
    ]
    
    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.level = 0

# --- Slide 3: Solution ---
def add_solution_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, content = slide.shapes.title, slide.placeholders[1]
    
    title.text = "Meet Your Digital Financial Twin"
    title.text_frame.paragraphs[0].font.color.rgb = IDBI_BLUE
    
    tf = content.text_frame
    tf.clear()
    
    features = [
        "Predictive Analytics: Loan eligibility, fraud risk, retirement planning.",
        "Scenario Simulator: 'What if I buy a car?' → AI simulates impact.",
        "Real-Time Insights: Cash flow forecasting, fraud alerts, investment readiness.",
        "Natural Language Queries: Ask 'Can I afford a home loan?'"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.level = 0

    # Add dashboard screenshot (placeholder)
    img_path = os.path.join(os.path.dirname(__file__), "dashboard_screenshot.png")
    if os.path.exists(img_path):
        left = Inches(6.5)
        top = Inches(2)
        slide.shapes.add_picture(img_path, left, top, width=Inches(6))

# --- Slide 4: Tech Stack ---
def add_tech_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, content = slide.shapes.title, slide.placeholders[1]
    
    title.text = "Tech Stack & Innovation"
    title.text_frame.paragraphs[0].font.color.rgb = IDBI_BLUE
    
    tf = content.text_frame
    tf.clear()
    
    tech_stack = [
        "Frontend: React + Tailwind + Plotly",
        "Backend: FastAPI + PostgreSQL + Redis",
        "ML Models: XGBoost (loan eligibility), CatBoost (fraud detection)",
        "Explainability: SHAP for RBI/IDBI compliance",
        "Innovation: Federated learning (privacy-preserving), NLP queries"
    ]
    
    for tech in tech_stack:
        p = tf.add_paragraph()
        p.text = tech
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.level = 0

# --- Slide 5: Business Impact ---
def add_impact_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, content = slide.shapes.title, slide.placeholders[1]
    
    title.text = "Business Impact & Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = IDBI_BLUE
    
    tf = content.text_frame
    tf.clear()
    
    impact = [
        "Banks: Reduce loan defaults by 30% with predictive scoring.",
        "Users: Save ₹50,000/year with personalized advice.",
        "Society: Financial inclusion for 500M+ Indians."
    ]
    
    for point in impact:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Roadmap:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IDBI_BLUE
    p.level = 0
    
    roadmap = [
        "Phase 1: Pilot with IDBI Bank (Q4 2026)",
        "Phase 2: Integrate with UPI/Aadhaar (2027)",
        "Phase 3: Expand to insurance/investments (2028)"
    ]
    
    for phase in roadmap:
        p = tf.add_paragraph()
        p.text = phase
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.level = 1

# --- Generate Slides ---
if __name__ == "__main__":
    print("Generating pitch deck...")
    add_title_slide()
    add_problem_slide()
    add_solution_slide()
    add_tech_slide()
    add_impact_slide()
    
    output_path = os.path.join(os.path.dirname(__file__), "slides.pptx")
    prs.save(output_path)
    print(f"[OK] Pitch deck saved: {output_path}")